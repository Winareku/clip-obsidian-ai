"""
extractors/__init__.py
----------------------
Text extraction layer.

Exports:
  - BaseExtractor      (ABC)
  - PDFExtractor
  - PPTXExtractor
  - DocxExtractor
  - TXTExtractor
  - ExtractorFactory   (creates the right extractor for a given file path)
"""

from __future__ import annotations

import logging
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


# ══════════════════════════════════════════════════════════════════════════════
# Abstract Base
# ══════════════════════════════════════════════════════════════════════════════

class BaseExtractor(ABC):
    """
    Interface for all text extractors.
    Every concrete extractor must implement `extract_text`.
    """

    @abstractmethod
    def extract_text(self, file_path: str) -> str:
        """
        Extract and return the plain-text content of *file_path*.
        Returns an empty string if extraction fails (never raises).
        """
        ...

    def _safe_extract(self, file_path: str) -> str:
        """
        Wraps extract_text with top-level error handling.
        Subclasses should override `extract_text`, not this method.
        """
        try:
            return self.extract_text(file_path)
        except Exception as exc:  # noqa: BLE001
            logger.error(
                "[%s] Failed to extract '%s': %s",
                self.__class__.__name__,
                file_path,
                exc,
            )
            return ""


# ══════════════════════════════════════════════════════════════════════════════
# Concrete Extractors
# ══════════════════════════════════════════════════════════════════════════════

class PDFExtractor(BaseExtractor):
    """Extracts text from PDF files using pdfplumber (with pypdf fallback)."""

    def extract_text(self, file_path: str) -> str:
        # Primary: pdfplumber (better layout preservation)
        try:
            import pdfplumber  # type: ignore

            text_parts: list[str] = []
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, start=1):
                    page_text = page.extract_text() or ""
                    if page_text.strip():
                        text_parts.append(f"[Página {page_num}]\n{page_text}")
            return "\n\n".join(text_parts)

        except ImportError:
            logger.debug("pdfplumber not available, falling back to pypdf.")

        # Fallback: pypdf
        try:
            from pypdf import PdfReader  # type: ignore

            reader = PdfReader(file_path)
            parts: list[str] = []
            for i, page in enumerate(reader.pages, start=1):
                text = page.extract_text() or ""
                if text.strip():
                    parts.append(f"[Página {i}]\n{text}")
            return "\n\n".join(parts)

        except ImportError as exc:
            raise ImportError(
                "Neither pdfplumber nor pypdf is installed. "
                "Run: pip install pdfplumber"
            ) from exc


class PPTXExtractor(BaseExtractor):
    """Extracts text from PowerPoint (.pptx) files using python-pptx."""

    def extract_text(self, file_path: str) -> str:
        try:
            from pptx import Presentation  # type: ignore
            from pptx.util import Pt  # type: ignore  # noqa: F401
        except ImportError as exc:
            raise ImportError(
                "python-pptx is not installed. Run: pip install python-pptx"
            ) from exc

        prs = Presentation(file_path)
        slide_texts: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            parts: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)

            if parts:
                slide_texts.append(f"[Diapositiva {slide_num}]\n" + "\n".join(parts))

        return "\n\n".join(slide_texts)


class DocxExtractor(BaseExtractor):
    """Extracts text from Word (.docx) files using python-docx."""

    def extract_text(self, file_path: str) -> str:
        try:
            import docx  # type: ignore
        except ImportError as exc:
            raise ImportError(
                "python-docx is not installed. Run: pip install python-docx"
            ) from exc

        doc = docx.Document(file_path)
        paragraphs: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]

        # Also extract text from tables
        table_texts: list[str] = []
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    table_texts.append(row_text)

        all_parts = paragraphs
        if table_texts:
            all_parts += ["\n[Tablas]\n"] + table_texts

        return "\n".join(all_parts)


class TXTExtractor(BaseExtractor):
    """Extracts text from plain text files (.txt, .md, etc.)."""

    # Encodings to try, in order of preference
    _ENCODINGS = ("utf-8", "utf-8-sig", "latin-1", "cp1252")

    def extract_text(self, file_path: str) -> str:
        for encoding in self._ENCODINGS:
            try:
                return Path(file_path).read_text(encoding=encoding)
            except UnicodeDecodeError:
                continue
        # Last resort: replace undecodable bytes
        return Path(file_path).read_text(encoding="utf-8", errors="replace")


# ══════════════════════════════════════════════════════════════════════════════
# Factory
# ══════════════════════════════════════════════════════════════════════════════

class ExtractorFactory:
    """
    Factory that maps file extensions to extractor classes.
    
    Usage:
        text = ExtractorFactory.extract("path/to/file.pdf")
    """

    _REGISTRY: dict[str, type[BaseExtractor]] = {
        ".pdf":  PDFExtractor,
        ".pptx": PPTXExtractor,
        ".docx": DocxExtractor,
        ".txt":  TXTExtractor,
        ".md":   TXTExtractor,
    }

    @classmethod
    def get_extractor(cls, file_path: str) -> Optional[BaseExtractor]:
        """Returns the appropriate extractor instance, or None if unsupported."""
        ext = Path(file_path).suffix.lower()
        extractor_cls = cls._REGISTRY.get(ext)
        if extractor_cls is None:
            logger.warning("No extractor registered for extension '%s'.", ext)
            return None
        return extractor_cls()

    @classmethod
    def extract(cls, file_path: str) -> str:
        """
        Convenience method: resolves the extractor and returns the text.
        Returns an empty string for unsupported or failed files.
        """
        extractor = cls.get_extractor(file_path)
        if extractor is None:
            return ""
        return extractor._safe_extract(file_path)

    @classmethod
    def supported_extensions(cls) -> list[str]:
        """Returns a sorted list of all registered extensions."""
        return sorted(cls._REGISTRY.keys())

    @classmethod
    def register(cls, extension: str, extractor_cls: type[BaseExtractor]) -> None:
        """
        Allows runtime registration of new extractors.
        Makes the system open for extension (Open/Closed Principle).
        
        Example:
            ExtractorFactory.register(".rst", TXTExtractor)
        """
        cls._REGISTRY[extension.lower()] = extractor_cls
        logger.debug("Registered extractor %s for '%s'.", extractor_cls.__name__, extension)
